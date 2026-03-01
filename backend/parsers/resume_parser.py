"""Resume parser supporting multiple formats: PDF, DOCX, PPTX, Images"""
import os
import logging
from typing import Dict, Any, Optional, Tuple
from pathlib import Path
import re
from datetime import datetime

from docx import Document
from pypdf import PdfReader
from pdf2image import convert_from_bytes
from pptx import Presentation

# Note: pytesseract requires system-level Tesseract installation
try:
    import pytesseract
    HAS_PYTESSERACT = True
except ImportError:
    HAS_PYTESSERACT = False

logger = logging.getLogger(__name__)


class ResumeParser:
    """Parser for multiple resume formats"""
    
    # Common section headers in resumes
    SECTION_PATTERNS = {
        'skills': r'(?:skills?|competencies|technical skills|expertise)',
        'experience': r'(?:experience|work experience|professional experience|employment)',
        'education': r'(?:education|academic|qualification)',
        'summary': r'(?:summary|objective|professional summary)',
        'contact': r'(?:contact|phone|email|linkedin)',
        'certifications': r'(?:certification|certificate|licenses?)',
    }
    
    def __init__(self):
        self.text_content = ""
        self.metadata = {}
    
    def parse_file(self, file_path: str) -> Dict[str, Any]:
        """
        Parse resume from file
        
        Args:
            file_path: Path to resume file
            
        Returns:
            Parsed resume data with metadata
        """
        file_ext = Path(file_path).suffix.lower()
        
        try:
            if file_ext == '.pdf':
                self.text_content = self._parse_pdf(file_path)
            elif file_ext == '.docx':
                self.text_content = self._parse_docx(file_path)
            elif file_ext == '.pptx':
                self.text_content = self._parse_pptx(file_path)
            elif file_ext in ['.jpg', '.jpeg', '.png', '.bmp']:
                self.text_content = self._parse_image(file_path)
            elif file_ext == '.txt':
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as _f:
                    self.text_content = _f.read()
            else:
                raise ValueError(f"Unsupported file format: {file_ext}")
            
            # Extract structured data from parsed text
            self.metadata = self._extract_metadata()
            
            return {
                'raw_text': self.text_content,
                'metadata': self.metadata,
                'parsed_at': datetime.now().isoformat(),
                'file_name': Path(file_path).name,
            }
        except Exception as e:
            logger.error(f"Error parsing resume {file_path}: {str(e)}")
            raise
    
    def _parse_pdf(self, file_path: str) -> str:
        """Parse PDF resume"""
        logger.info(f"Parsing PDF: {file_path}")
        text = ""
        
        try:
            with open(file_path, 'rb') as f:
                pdf_reader = PdfReader(f)
                for page_num, page in enumerate(pdf_reader.pages):
                    text += page.extract_text()
        except Exception as e:
            logger.warning(f"Text extraction failed for PDF, trying OCR: {e}")
            # Fallback to OCR for scanned PDFs
            if HAS_PYTESSERACT:
                text = self._ocr_pdf(file_path)
            else:
                logger.error("Pytesseract not available for OCR fallback")
                raise
        
        return text
    
    def _ocr_pdf(self, file_path: str) -> str:
        """Use OCR on PDF if text extraction fails"""
        try:
            with open(file_path, 'rb') as f:
                images = convert_from_bytes(f.read())
            
            text = ""
            for image in images:
                text += pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.error(f"OCR failed for PDF: {e}")
            raise
    
    def _parse_docx(self, file_path: str) -> str:
        """Parse DOCX (Microsoft Word) resume"""
        logger.info(f"Parsing DOCX: {file_path}")
        doc = Document(file_path)
        text = ""
        
        for paragraph in doc.paragraphs:
            text += paragraph.text + "\n"
        
        # Extract text from tables
        for table in doc.tables:
            for row in table.rows:
                for cell in row.cells:
                    text += cell.text + " "
                text += "\n"
        
        return text
    
    def _parse_pptx(self, file_path: str) -> str:
        """Parse PPTX (PowerPoint) resume"""
        logger.info(f"Parsing PPTX: {file_path}")
        prs = Presentation(file_path)
        text = ""
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        
        return text
    
    def _parse_image(self, file_path: str) -> str:
        """Parse image resume using OCR"""
        if not HAS_PYTESSERACT:
            raise RuntimeError("Pytesseract required for image parsing. Install: pip install pytesseract")
        
        logger.info(f"Parsing image resume with OCR: {file_path}")
        try:
            from PIL import Image
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            return text
        except Exception as e:
            logger.error(f"OCR failed for image: {e}")
            raise
    
    def _extract_metadata(self) -> Dict[str, Any]:
        """Extract structured metadata from resume text"""
        metadata = {
            'candidate_name': self._extract_name(),
            'email': self._extract_email(),
            'phone': self._extract_phone(),
            'skills': self._extract_skills(),
            'years_of_experience': self._estimate_experience_years(),
            'education_level': self._extract_education_level(),
            'current_role': self._extract_current_role(),
        }
        
        return {k: v for k, v in metadata.items() if v is not None}
    
    def _extract_name(self) -> Optional[str]:
        """Extract candidate name from resume"""
        # Simple heuristic: first line often contains name
        lines = self.text_content.strip().split('\n')
        if lines:
            # Filter out very short lines and common section headers
            for line in lines[:10]:  # Check first 10 lines
                cleaned = line.strip()
                if (len(cleaned) > 5 and 
                    len(cleaned) < 100 and 
                    not any(keyword in cleaned.lower() for keyword in ['resume', 'cv', 'mailto', 'linkedin'])):
                    return cleaned
        return None
    
    def _extract_email(self) -> Optional[str]:
        """Extract email address from resume"""
        email_pattern = r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b'
        match = re.search(email_pattern, self.text_content)
        return match.group(0) if match else None
    
    def _extract_phone(self) -> Optional[str]:
        """Extract phone number from resume"""
        # Common phone number patterns
        phone_patterns = [
            r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b',  # XXX-XXX-XXXX
            r'\b\(\d{3}\)\s?\d{3}[-.]?\d{4}\b',  # (XXX) XXX-XXXX
            r'\b\+\d{1,3}\s?\d{1,14}\b',  # International format
        ]
        
        for pattern in phone_patterns:
            match = re.search(pattern, self.text_content)
            if match:
                return match.group(0)
        return None
    
    def _extract_skills(self) -> list:
        """Extract technical skills from resume"""
        # Common technical skills to look for
        common_skills = [
            'python', 'java', 'javascript', 'c\\+\\+', 'c#', 'ruby', 'go', 'rust', 'php',
            'sql', 'mysql', 'postgresql', 'mongodb', 'redis',
            'react', 'angular', 'vue', 'vue.js', 'django', 'flask', 'fastapi',
            'aws', 'azure', 'gcp', 'docker', 'kubernetes',
            'git', 'jenkins', 'ci/cd', 'devops',
            'machine learning', 'deep learning', 'tensorflow', 'pytorch', 'scikit-learn',
            'api', 'rest', 'graphql', 'html', 'css', 'scss',
            'agile', 'scrum', 'jira', 'confluence',
            'linux', 'windows', 'unix',
            'excel', 'power bi', 'tableau', 'looker',
        ]
        
        text_lower = self.text_content.lower()
        found_skills = []
        
        for skill in common_skills:
            # Use word boundaries to avoid partial matches
            pattern = r'\b' + skill.replace('+', r'\+').replace('/', r'\/') + r'\b'
            if re.search(pattern, text_lower):
                found_skills.append(skill.replace('\\', ''))
        
        return found_skills
    
    def _estimate_experience_years(self) -> Optional[int]:
        """Estimate years of experience from resume"""
        # Look for year patterns like "2010 - 2020" or "2010-2020"
        year_pattern = r'(20\d{2})\s*[-–]\s*(20\d{2}|present)'
        matches = re.findall(year_pattern, self.text_content, re.IGNORECASE)
        
        if matches:
            total_years = 0
            for start, end in matches:
                if 'present' in end.lower():
                    end = str(datetime.now().year)
                total_years += int(end) - int(start)
            return total_years if total_years > 0 else None
        
        return None
    
    def _extract_education_level(self) -> Optional[str]:
        """Extract education level from resume"""
        education_keywords = {
            'phd': ['phd', 'doctorate', 'dr.', 'ph.d'],
            'master': ['master', 'ms', 'mba', 'm.s', 'ma', 'm.a'],
            'bachelor': ['bachelor', 'bs', 'ba', 'b.s', 'b.a', 'btech'],
            'associate': ['associate', 'diploma', 'associate degree'],
            'high_school': ['high school', 'hs', 'secondary'],
        }
        
        text_lower = self.text_content.lower()
        
        for level, keywords in education_keywords.items():
            for keyword in keywords:
                if keyword in text_lower:
                    return level
        
        return None
    
    def _extract_current_role(self) -> Optional[str]:
        """Extract current job title from resume"""
        # Look for job titles near the top of the document
        text_lines = self.text_content.split('\n')[:20]
        
        job_keywords = [
            'engineer', 'developer', 'manager', 'lead', 'architect', 'analyst',
            'consultant', 'specialist', 'coordinator', 'director', 'officer',
            'associate', 'senior', 'junior', 'principal', 'staff',
        ]
        
        for line in text_lines:
            if any(keyword in line.lower() for keyword in job_keywords):
                return line.strip()
        
        return None
